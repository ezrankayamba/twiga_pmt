from pptx import Presentation
from . import models


def import_n_replace(file):
    print('Name: ', file.name)
    prs = Presentation(file)
    text_runs = []
    for slide in prs.slides:
        slide_runs = []
        title = 'No set'
        winner = 'Not set'
        facts = []
        # print()
        for shape in slide.shapes:
            # print('Shape: ', shape.name, shape.shape_type)
            if shape.name == 'Title 1':
                title = shape.text.strip()
            if shape.name == 'Content Placeholder 2':
                parts = list(filter(lambda x: x, shape.text.splitlines()))
                winner = parts[0]
                facts = parts[1:]
        text_runs.append({'title': title, 'winner': winner, 'facts': facts})

    # Clear DB
    models.AwardSession.objects.all().delete()

    session = models.AwardSession(title=text_runs[0]['title'])
    session.save()

    for run in text_runs[1:]:
        win = models.Win(title=run['title'], winner=run['winner'], award_session=session)
        win.save()

        for f in run['facts']:
            fact = models.Fact(name=f, win=win)
            fact.save()
